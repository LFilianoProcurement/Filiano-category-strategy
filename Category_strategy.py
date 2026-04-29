import streamlit as st
import anthropic
import json
from datetime import datetime
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker


# ── PPTX Generator ────────────────────────────────────────────────────────────
def hex_color(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

DARK  = "0F2027"
MID   = "2C5364"
LIGHT = "EBF4F8"
WHITE = "FFFFFF"
ACCENT = "00B4D8"
GRAY  = "4A5568"
LGRAY = "E8EDF2"

def add_bg(slide, color=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_color(color)

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_color(color)
    shape.line.fill.background()
    return shape

def add_frame(slide, l, t, w, h, border_color=MID, border_width=1.5):
    """Rectangle with no fill, just a visible border frame."""
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.background()
    shape.line.color.rgb = hex_color(border_color)
    shape.line.width = Pt(border_width)
    return shape

def add_text(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = hex_color(color)
    run.font.name = "Calibri"
    return txBox


# ── Chart helpers ─────────────────────────────────────────────────────────────
CHART_COLORS = ["#0F2027", "#2C5364", "#00B4D8", "#48CAE4", "#90E0EF",
                "#ADE8F4", "#CAF0F8", "#023E8A", "#0077B6", "#0096C7"]

def detect_columns(df):
    """Auto-detect supplier, period, and amount columns by name heuristics."""
    cols = {c.lower(): c for c in df.columns}
    supplier_col = next((cols[c] for c in cols if any(k in c for k in
                        ["supplier","vendor","company","provider","name"])), df.columns[0])
    amount_col   = next((cols[c] for c in cols if any(k in c for k in
                        ["amount","spend","cost","value","total","usd","$"])), df.columns[-1])
    period_col   = next((cols[c] for c in cols if any(k in c for k in
                        ["period","date","quarter","month","year","time","q1","q2","q3","q4"])), None)
    # Force amount column to numeric, stripping any $ or commas
    df[amount_col] = pd.to_numeric(
        df[amount_col].astype(str).str.replace(r"[\$,]", "", regex=True).str.strip(),
        errors="coerce"
    ).fillna(0)
    return supplier_col, amount_col, period_col

def fig_to_image(fig):
    """Convert matplotlib figure to bytes buffer."""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf

def chart_supplier_bar(df, supplier_col, amount_col):
    """Horizontal bar: spend by supplier."""
    data = df.groupby(supplier_col)[amount_col].sum().sort_values()
    fig, ax = plt.subplots(figsize=(7, max(3, len(data) * 0.55)))
    fig.patch.set_facecolor("white")
    bars = ax.barh(data.index, data.values,
                   color=CHART_COLORS[:len(data)], edgecolor="none", height=0.6)
    ax.set_facecolor("white")
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(
        lambda x, _: f"${x/1_000_000:.1f}M" if x >= 1_000_000 else f"${x/1_000:.0f}K"))
    ax.spines[["top","right","left"]].set_visible(False)
    ax.tick_params(axis="y", labelsize=9, colors="#2d3748")
    ax.tick_params(axis="x", labelsize=8, colors="#718096")
    ax.set_title("Spend by Supplier", fontsize=12, fontweight="bold",
                 color="#0F2027", pad=10, loc="left")
    for bar, val in zip(bars, data.values):
        label = f"${val/1_000_000:.2f}M" if val >= 1_000_000 else f"${val/1_000:.0f}K"
        ax.text(val + data.max() * 0.01, bar.get_y() + bar.get_height() / 2,
                label, va="center", fontsize=8, color="#2d3748")
    fig.tight_layout()
    return fig_to_image(fig)

def chart_spend_trend(df, period_col, amount_col):
    """Line chart: spend over time."""
    data = df.groupby(period_col)[amount_col].sum()
    fig, ax = plt.subplots(figsize=(7, 3.5))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    ax.plot(data.index.astype(str), data.values, color="#00B4D8",
            linewidth=2.5, marker="o", markersize=6, markerfacecolor="#0F2027")
    ax.fill_between(range(len(data)), data.values, alpha=0.08, color="#00B4D8")
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(
        lambda x, _: f"${x/1_000_000:.1f}M" if x >= 1_000_000 else f"${x/1_000:.0f}K"))
    ax.spines[["top","right"]].set_visible(False)
    ax.tick_params(axis="x", labelsize=8, rotation=30, colors="#718096")
    ax.tick_params(axis="y", labelsize=8, colors="#718096")
    ax.set_title("Spend Over Time", fontsize=12, fontweight="bold",
                 color="#0F2027", pad=10, loc="left")
    fig.tight_layout()
    return fig_to_image(fig)

def chart_supplier_pie(df, supplier_col, amount_col):
    """Donut chart: supplier share."""
    data = df.groupby(supplier_col)[amount_col].sum().sort_values(ascending=False)
    if len(data) > 6:
        top = data.iloc[:5]
        top["Other"] = data.iloc[5:].sum()
        data = top
    fig, ax = plt.subplots(figsize=(5, 4))
    fig.patch.set_facecolor("white")
    wedges, texts, autotexts = ax.pie(
        data.values, labels=data.index,
        colors=CHART_COLORS[:len(data)],
        autopct="%1.1f%%", startangle=90,
        wedgeprops={"width": 0.55, "edgecolor": "white", "linewidth": 2},
        pctdistance=0.75
    )
    for t in texts: t.set_fontsize(8)
    for a in autotexts: a.set_fontsize(7); a.set_color("white")
    ax.set_title("Supplier Concentration", fontsize=12, fontweight="bold",
                 color="#0F2027", pad=10, loc="left")
    fig.tight_layout()
    return fig_to_image(fig)

def chart_top_suppliers_trend(df, supplier_col, period_col, amount_col):
    """Multi-line: top 5 suppliers over time."""
    top5 = df.groupby(supplier_col)[amount_col].sum().nlargest(5).index
    data = df[df[supplier_col].isin(top5)].groupby(
        [period_col, supplier_col])[amount_col].sum().unstack(fill_value=0)
    fig, ax = plt.subplots(figsize=(7, 3.5))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    for i, col in enumerate(data.columns):
        ax.plot(data.index.astype(str), data[col],
                color=CHART_COLORS[i], linewidth=2, marker="o",
                markersize=5, label=col)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(
        lambda x, _: f"${x/1_000_000:.1f}M" if x >= 1_000_000 else f"${x/1_000:.0f}K"))
    ax.spines[["top","right"]].set_visible(False)
    ax.tick_params(axis="x", labelsize=7, rotation=30, colors="#718096")
    ax.tick_params(axis="y", labelsize=7, colors="#718096")
    ax.legend(fontsize=7, frameon=False)
    ax.set_title("Top Suppliers Over Time", fontsize=12, fontweight="bold",
                 color="#0F2027", pad=10, loc="left")
    fig.tight_layout()
    return fig_to_image(fig)


def add_spend_slide(prs, blank, df, cat):
    """Insert a Spend Analysis slide and return the slide."""
    supplier_col, amount_col, period_col = detect_columns(df)
    total = df[amount_col].sum()

    s = prs.slides.add_slide(blank)
    add_bg(s, "FFFFFF")
    add_frame(s, 0.3, 0.15, 12.7, 0.75)
    add_rect(s, 0.3, 0.15, 0.12, 0.75, ACCENT)
    add_text(s, "SPEND ANALYSIS", 0.6, 0.22, 10, 0.6, size=18, bold=True, color=DARK)

    # Total spend callout
    total_str = f"${total/1_000_000:.2f}M" if total >= 1_000_000 else f"${total/1_000:.0f}K"
    num_suppliers = df[supplier_col].nunique()
    add_frame(s, 0.3, 1.05, 3.0, 0.9)
    add_text(s, total_str, 0.45, 1.1, 2.7, 0.5, size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, "Total Analyzed Spend", 0.45, 1.55, 2.7, 0.3, size=8, color=GRAY, align=PP_ALIGN.CENTER)
    add_frame(s, 3.45, 1.05, 2.5, 0.9)
    add_text(s, str(num_suppliers), 3.6, 1.1, 2.2, 0.5, size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, "Suppliers", 3.6, 1.55, 2.2, 0.3, size=8, color=GRAY, align=PP_ALIGN.CENTER)

    # Supplier bar chart (left)
    bar_img = chart_supplier_bar(df, supplier_col, amount_col)
    s.shapes.add_picture(bar_img, Inches(0.3), Inches(2.1), Inches(6.5), Inches(4.8))

    # Pie chart (right top)
    pie_img = chart_supplier_pie(df, supplier_col, amount_col)
    s.shapes.add_picture(pie_img, Inches(6.9), Inches(1.0), Inches(6.0), Inches(3.0))

    # Trend or multi-line chart (right bottom)
    if period_col:
        if df[period_col].nunique() > 1:
            trend_img = chart_spend_trend(df, period_col, amount_col)
            s.shapes.add_picture(trend_img, Inches(6.9), Inches(4.1), Inches(6.0), Inches(2.9))

    return s


def build_pptx(d, inp, spend_df=None):
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # completely blank

    cat = inp.get("category_name", "Category Strategy")
    spend = inp.get("annual_spend", 0)
    ov = d.get("category_overview", {})
    pf = d.get("porter_five_forces", {})
    sw = d.get("swot", {})
    ss = d.get("sourcing_strategy", {})
    sv = d.get("savings_opportunity", {})

    # ── Slide 1 – Title ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, DARK)
    add_rect(s, 0, 0, 13.33, 7.5, MID)
    add_rect(s, 0, 0, 0.18, 7.5, ACCENT)
    add_text(s, "CATEGORY STRATEGY", 0.5, 1.2, 12, 0.5, size=11, color=ACCENT, bold=True)
    add_text(s, cat, 0.5, 1.8, 12, 1.4, size=38, bold=True, color=WHITE)
    add_text(s, inp.get("industry", ""), 0.5, 3.3, 8, 0.5, size=16, color="A8C0CC")
    add_text(s, f"Annual Spend: ${spend:,}  |  Suppliers: {inp.get('num_suppliers','')}  |  Contract Coverage: {inp.get('contract_coverage','')}%",
             0.5, 4.0, 11, 0.4, size=13, color="A8C0CC")
    add_text(s, f"Generated {datetime.now().strftime('%B %d, %Y')}",
             0.5, 6.8, 6, 0.4, size=10, color="6B8A99")

    # ── Slide 1b – Spend Analysis (optional) ─────────────────────────────────
    if spend_df is not None:
        add_spend_slide(prs, blank, spend_df, cat)

    # ── Slide 2 – Executive Summary ───────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_frame(s, 0.3, 0.15, 12.7, 0.75)
    add_rect(s, 0.3, 0.15, 0.12, 0.75, ACCENT)
    add_text(s, "EXECUTIVE SUMMARY", 0.6, 0.22, 10, 0.6, size=18, bold=True, color=DARK)

    # Kraljic + savings boxes
    pos = ov.get("kraljic_position", "—")
    attract = pf.get("overall_attractiveness", "—")
    sav_pct = sv.get("estimated_savings_pct", "—")
    sav_usd = sv.get("estimated_savings_usd", "—")

    for i, (label, val) in enumerate([
        ("Kraljic Position", pos),
        ("Market Attractiveness", attract),
        ("Savings Range", sav_pct),
        ("Est. Savings Value", sav_usd),
    ]):
        x = 0.4 + i * 3.2
        add_frame(s, x, 1.15, 3.0, 1.1)
        add_text(s, val, x + 0.15, 1.25, 2.7, 0.55, size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(s, label, x + 0.15, 1.75, 2.7, 0.35, size=9, color=GRAY, align=PP_ALIGN.CENTER)

    add_text(s, d.get("executive_summary", ""), 0.4, 2.5, 12.5, 1.6, size=13, color=GRAY)
    add_text(s, "Market Outlook", 0.4, 4.2, 6, 0.35, size=11, bold=True, color=DARK)
    add_text(s, ov.get("market_outlook", ""), 0.4, 4.6, 6, 1.5, size=12, color=GRAY)
    add_text(s, "Kraljic Rationale", 6.8, 4.2, 6, 0.35, size=11, bold=True, color=DARK)
    add_text(s, ov.get("kraljic_rationale", ""), 6.8, 4.6, 6, 1.5, size=12, color=GRAY)

    # ── Slide 3 – Porter's Five Forces ────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_frame(s, 0.3, 0.15, 12.7, 0.75)
    add_rect(s, 0.3, 0.15, 0.12, 0.75, ACCENT)
    add_text(s, "PORTER'S FIVE FORCES", 0.6, 0.22, 10, 0.6, size=18, bold=True, color=DARK)

    forces = [
        ("Supplier Power",       "supplier_power"),
        ("Buyer Power",          "buyer_power"),
        ("Threat of Substitutes","threat_of_substitutes"),
        ("New Entrants",         "threat_of_new_entrants"),
        ("Competitive Rivalry",  "competitive_rivalry"),
    ]
    rating_colors = {"High": "C53030", "Medium": "92400E", "Low": "065F46"}
    rating_bg     = {"High": "FDE8E8", "Medium": "FEF3C7", "Low": "D1FAE5"}

    positions = [(0.3, 1.1), (4.4, 1.1), (8.5, 1.1), (2.35, 3.8), (6.45, 3.8)]
    for idx, (label, key) in enumerate(forces):
        force = pf.get(key, {})
        rating = force.get("rating", "Medium")
        analysis = force.get("analysis", "")
        x, y = positions[idx]
        add_frame(s, x, y, 4.3, 2.4)
        add_rect(s, x, y, 4.3, 0.42, WHITE)
        add_text(s, label, x + 0.15, y + 0.07, 2.8, 0.32, size=11, bold=True, color=DARK)
        add_rect(s, x + 3.1, y + 0.06, 1.0, 0.3, rating_bg.get(rating, "FEF3C7"))
        add_text(s, rating, x + 3.1, y + 0.06, 1.0, 0.3, size=9, bold=True,
                 color=rating_colors.get(rating, "92400E"), align=PP_ALIGN.CENTER)
        add_text(s, analysis, x + 0.15, y + 0.52, 3.95, 1.75, size=10, color=GRAY)

    attract = pf.get("overall_attractiveness", "Neutral")
    a_bg = {"Favorable": "D1FAE5", "Neutral": "FEF3C7", "Challenging": "FDE8E8"}.get(attract, "FEF3C7")
    a_tc = {"Favorable": "065F46", "Neutral": "92400E", "Challenging": "C53030"}.get(attract, "92400E")
    add_rect(s, 3.5, 6.65, 6.3, 0.55, a_bg)
    add_text(s, f"Overall Market Attractiveness: {attract}", 3.5, 6.65, 6.3, 0.55,
             size=12, bold=True, color=a_tc, align=PP_ALIGN.CENTER)

    # ── Slide 4 – SWOT ────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_frame(s, 0.3, 0.15, 12.7, 0.75)
    add_rect(s, 0.3, 0.15, 0.12, 0.75, ACCENT)
    add_text(s, "SWOT ANALYSIS", 0.6, 0.22, 10, 0.6, size=18, bold=True, color=DARK)

    swot_config = [
        ("STRENGTHS",    "strengths",    "D1FAE5", "065F46", 0.3,  1.1),
        ("WEAKNESSES",   "weaknesses",   "FDE8E8", "C53030", 6.85, 1.1),
        ("OPPORTUNITIES","opportunities","DBEAFE", "1E40AF", 0.3,  4.2),
        ("THREATS",      "threats",      "FEF3C7", "92400E", 6.85, 4.2),
    ]
    for title, key, bg, tc, x, y in swot_config:
        items = sw.get(key, [])
        add_rect(s, x, y, 6.25, 2.9, bg)
        add_text(s, title, x + 0.2, y + 0.15, 5.8, 0.4, size=12, bold=True, color=tc)
        for i, item in enumerate(items[:4]):
            add_text(s, f"• {item}", x + 0.2, y + 0.6 + i * 0.54, 5.8, 0.5, size=11, color="2D3748")

    # ── Slide 5 – Sourcing Strategy ───────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_frame(s, 0.3, 0.15, 12.7, 0.75)
    add_rect(s, 0.3, 0.15, 0.12, 0.75, ACCENT)
    add_text(s, "SOURCING STRATEGY", 0.6, 0.22, 10, 0.6, size=18, bold=True, color=DARK)

    add_frame(s, 0.3, 1.1, 12.7, 0.55)
    add_text(s, f"Recommended Approach: {ss.get('recommended_approach', '—')}",
             0.5, 1.15, 12, 0.45, size=14, bold=True, color=DARK)

    add_text(s, "Rationale", 0.4, 1.85, 6, 0.35, size=11, bold=True, color=DARK)
    add_text(s, ss.get("rationale", ""), 0.4, 2.25, 6, 1.6, size=12, color=GRAY)

    add_text(s, "Supplier Segmentation", 0.4, 4.0, 6, 0.35, size=11, bold=True, color=DARK)
    add_text(s, ss.get("supplier_segmentation", ""), 0.4, 4.4, 6, 1.7, size=12, color=GRAY)

    add_text(s, "Negotiation Levers", 7.0, 1.85, 5.9, 0.35, size=11, bold=True, color=DARK)
    levers = ss.get("negotiation_leverage", [])
    for i, lever in enumerate(levers[:5]):
        add_frame(s, 7.0, 2.3 + i * 0.7, 5.9, 0.55)
        add_text(s, f"🔧 {lever}", 7.15, 2.33 + i * 0.7, 5.6, 0.45, size=11, color=DARK)

    # ── Slide 6 – Key Initiatives ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_frame(s, 0.3, 0.15, 12.7, 0.75)
    add_rect(s, 0.3, 0.15, 0.12, 0.75, ACCENT)
    add_text(s, "KEY INITIATIVES", 0.6, 0.22, 10, 0.6, size=18, bold=True, color=DARK)

    initiatives = d.get("key_initiatives", [])
    p_colors = {"High": "C53030", "Medium": "92400E", "Low": "065F46"}
    p_bgs    = {"High": "FDE8E8", "Medium": "FEF3C7", "Low": "D1FAE5"}
    for i, ini in enumerate(initiatives[:5]):
        y = 1.1 + i * 1.25
        p = ini.get("priority", "Medium")
        add_rect(s, 0.3, y, 12.7, 1.1, LGRAY)
        add_rect(s, 0.3, y, 0.08, 1.1, p_colors.get(p, "92400E"))
        add_text(s, ini.get("initiative", ""), 0.55, y + 0.08, 7.5, 0.38, size=12, bold=True, color=DARK)
        add_text(s, ini.get("description", ""), 0.55, y + 0.5, 8.5, 0.5, size=10, color=GRAY)
        add_rect(s, 9.2, y + 0.08, 1.7, 0.3, p_bgs.get(p, "FEF3C7"))
        add_text(s, p, 9.2, y + 0.08, 1.7, 0.3, size=9, bold=True,
                 color=p_colors.get(p, "92400E"), align=PP_ALIGN.CENTER)
        add_text(s, ini.get("timeline", ""), 11.1, y + 0.08, 1.9, 0.3, size=9, color=MID)
        add_text(s, f"✓ {ini.get('expected_outcome', '')}", 0.55, y + 0.78, 12, 0.28,
                 size=9, color="2C5364", italic=True)

    # ── Slide 7 – Risk Register ───────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_frame(s, 0.3, 0.15, 12.7, 0.75)
    add_rect(s, 0.3, 0.15, 0.12, 0.75, ACCENT)
    add_text(s, "RISK REGISTER", 0.6, 0.22, 10, 0.6, size=18, bold=True, color=DARK)

    headers = ["Risk", "Likelihood", "Impact", "Mitigation"]
    col_x   = [0.3, 5.5, 7.3, 9.1]
    col_w   = [5.0, 1.6, 1.6, 4.1]
    add_frame(s, 0.3, 1.1, 12.9, 0.45)
    for hdr, x, w in zip(headers, col_x, col_w):
        add_text(s, hdr, x + 0.1, 1.13, w, 0.38, size=11, bold=True, color=DARK)

    risks = d.get("risk_register", [])
    for i, risk in enumerate(risks[:5]):
        y = 1.65 + i * 1.0
        add_frame(s, 0.3, y, 12.9, 0.9)
        add_text(s, risk.get("risk", ""), col_x[0] + 0.1, y + 0.05, col_w[0], 0.8, size=10, color=DARK)
        lh = risk.get("likelihood", "Medium")
        im = risk.get("impact", "Medium")
        for val, xi, wi in [(lh, col_x[1], col_w[1]), (im, col_x[2], col_w[2])]:
            rc = {"High": "C53030", "Medium": "92400E", "Low": "065F46"}.get(val, "92400E")
            rb = {"High": "FDE8E8", "Medium": "FEF3C7", "Low": "D1FAE5"}.get(val, "FEF3C7")
            add_rect(s, xi + 0.1, y + 0.2, wi - 0.2, 0.38, rb)
            add_text(s, val, xi + 0.1, y + 0.2, wi - 0.2, 0.38, size=9, bold=True,
                     color=rc, align=PP_ALIGN.CENTER)
        add_text(s, risk.get("mitigation", ""), col_x[3] + 0.1, y + 0.05, col_w[3], 0.8, size=10, color=GRAY)

    # ── Slide 8 – KPIs ────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_frame(s, 0.3, 0.15, 12.7, 0.75)
    add_rect(s, 0.3, 0.15, 0.12, 0.75, ACCENT)
    add_text(s, "KPIs & PERFORMANCE METRICS", 0.6, 0.22, 10, 0.6, size=18, bold=True, color=DARK)

    kpis = d.get("kpis", [])
    for i, kpi in enumerate(kpis[:5]):
        x = 0.3 + (i % 3) * 4.35
        y = 1.2 if i < 3 else 4.0
        add_frame(s, x, y, 4.1, 2.3)
        add_rect(s, x, y, 4.1, 0.08, ACCENT)
        add_text(s, kpi.get("metric", ""), x + 0.2, y + 0.2, 3.7, 0.8,
                 size=13, bold=True, color=DARK)
        add_text(s, kpi.get("target", ""), x + 0.2, y + 1.05, 3.7, 0.65,
                 size=12, color=ACCENT)
        add_text(s, kpi.get("frequency", ""), x + 0.2, y + 1.75, 3.7, 0.35,
                 size=9, color=GRAY)

    # ── Slide 9 – Closing ─────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, WHITE)
    add_rect(s, 0, 0, 0.18, 7.5, ACCENT)
    add_rect(s, 0, 2.8, 13.33, 0.08, ACCENT)
    add_text(s, cat, 0.5, 1.0, 12, 1.2, size=36, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, "Category Strategy Summary", 0.5, 2.2, 12, 0.5, size=14,
             color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, f"Savings Opportunity: {sv.get('estimated_savings_pct','—')}  |  Primary Lever: {sv.get('primary_lever','—')}  |  Timeframe: {sv.get('timeframe','—')}",
             0.5, 3.2, 12, 0.5, size=13, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, f"Prepared using AI-Powered Category Strategy Generator  |  {datetime.now().strftime('%B %Y')}",
             0.5, 6.8, 12, 0.4, size=9, color=GRAY, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

# ── Page config ──────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Category Strategy Generator",
    page_icon="🎯",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── Custom CSS ────────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=DM+Serif+Display:ital@0;1&family=DM+Sans:wght@300;400;500;600&display=swap');

html, body, [class*="css"] {
    font-family: 'DM Sans', sans-serif;
}

.main-header {
    background: linear-gradient(135deg, #0f2027 0%, #203a43 50%, #2c5364 100%);
    padding: 2.5rem 2rem;
    border-radius: 12px;
    margin-bottom: 2rem;
    color: white;
}

.main-header h1 {
    font-family: 'DM Serif Display', serif;
    font-size: 2.4rem;
    margin: 0 0 0.4rem 0;
    letter-spacing: -0.5px;
}

.main-header p {
    margin: 0;
    color: #a8c0cc;
    font-size: 1rem;
    font-weight: 300;
}

.section-card {
    background: #ffffff;
    border: 1px solid #e8edf2;
    border-radius: 10px;
    padding: 1.5rem;
    margin-bottom: 1.2rem;
    box-shadow: 0 2px 8px rgba(0,0,0,0.04);
}

.section-label {
    font-family: 'DM Serif Display', serif;
    font-size: 1.15rem;
    color: #0f2027;
    margin-bottom: 0.6rem;
    display: flex;
    align-items: center;
    gap: 0.5rem;
}

.tag-pill {
    display: inline-block;
    background: #e8f4f8;
    color: #2c5364;
    padding: 0.2rem 0.7rem;
    border-radius: 20px;
    font-size: 0.78rem;
    font-weight: 500;
    margin: 0.2rem;
}

.metric-box {
    background: linear-gradient(135deg, #0f2027, #2c5364);
    color: white;
    border-radius: 10px;
    padding: 1.2rem;
    text-align: center;
}

.metric-box .value {
    font-family: 'DM Serif Display', serif;
    font-size: 1.8rem;
}

.metric-box .label {
    font-size: 0.78rem;
    color: #a8c0cc;
    text-transform: uppercase;
    letter-spacing: 0.5px;
}

.stTabs [data-baseweb="tab-list"] {
    gap: 6px;
    background: #f0f4f8;
    padding: 6px;
    border-radius: 10px;
}

.stTabs [data-baseweb="tab"] {
    border-radius: 8px;
    font-size: 0.85rem;
    font-weight: 500;
    color: #4a5568;
    padding: 0.4rem 1rem;
}

.stTabs [aria-selected="true"] {
    background: white !important;
    color: #0f2027 !important;
    box-shadow: 0 2px 6px rgba(0,0,0,0.08);
}

.strategy-output {
    background: #f8fafc;
    border-left: 4px solid #2c5364;
    padding: 1.4rem 1.6rem;
    border-radius: 0 10px 10px 0;
    margin: 0.8rem 0;
    line-height: 1.7;
    font-size: 0.92rem;
}

.porter-grid {
    display: grid;
    grid-template-columns: 1fr 1fr;
    gap: 0.8rem;
    margin-top: 0.6rem;
}

.porter-card {
    background: white;
    border: 1px solid #dde4ec;
    border-radius: 8px;
    padding: 1rem;
}

.porter-card .force-title {
    font-weight: 600;
    font-size: 0.85rem;
    color: #0f2027;
    margin-bottom: 0.4rem;
    text-transform: uppercase;
    letter-spacing: 0.4px;
}

.rating-badge {
    display: inline-block;
    padding: 0.15rem 0.6rem;
    border-radius: 12px;
    font-size: 0.75rem;
    font-weight: 600;
    margin-left: 0.5rem;
}

.rating-high { background: #fde8e8; color: #c53030; }
.rating-medium { background: #fef3c7; color: #92400e; }
.rating-low { background: #d1fae5; color: #065f46; }

.download-btn {
    background: linear-gradient(135deg, #0f2027, #2c5364);
    color: white;
    border: none;
    padding: 0.6rem 1.5rem;
    border-radius: 8px;
    font-weight: 500;
    cursor: pointer;
}

.stButton > button {
    background: linear-gradient(135deg, #0f2027 0%, #2c5364 100%);
    color: white;
    border: none;
    border-radius: 8px;
    font-weight: 500;
    padding: 0.6rem 2rem;
    font-family: 'DM Sans', sans-serif;
    transition: opacity 0.2s;
}

.stButton > button:hover {
    opacity: 0.9;
    color: white;
    border: none;
}

.sidebar-label {
    font-size: 0.8rem;
    font-weight: 600;
    color: #4a5568;
    text-transform: uppercase;
    letter-spacing: 0.5px;
    margin-bottom: 0.3rem;
}

div[data-testid="stSidebar"] {
    background: #0f2027;
}

div[data-testid="stSidebar"] label,
div[data-testid="stSidebar"] .stMarkdown {
    color: #cbd5e0 !important;
}

div[data-testid="stSidebar"] h3 {
    color: white !important;
    font-family: 'DM Serif Display', serif;
}
</style>
""", unsafe_allow_html=True)

# ── Auth state ────────────────────────────────────────────────────────────────
if "authenticated" not in st.session_state:
    st.session_state.authenticated = False

# ── Header ────────────────────────────────────────────────────────────────────
st.markdown("""
<div class="main-header">
    <h1>🎯 Category Strategy Generator</h1>
    <p>AI-powered procurement strategy · Porter's Five Forces · Kraljic Positioning · Sourcing Roadmap</p>
</div>
""", unsafe_allow_html=True)

# ── Sidebar Inputs ─────────────────────────────────────────────────────────────
with st.sidebar:
    st.markdown("### ⚙️ Category Inputs")
    st.markdown("---")

    # ── Basic Info ────────────────────────────────────────────────────────────
    st.markdown("**📋 Basic Information**")
    category_name = st.text_input("Category Name", placeholder="e.g. ETO Sterilization Services")
    industry = st.selectbox("Industry / Sector", [
        "Medical Devices", "Pharmaceutical", "Life Sciences", "Healthcare",
        "Aerospace & Defense", "Automotive", "Chemical", "Oil & Gas",
        "Financial Services", "Technology", "Consumer Goods", "Other"
    ])
    annual_spend = st.number_input("Annual Spend (USD)", min_value=0, value=5_000_000, step=500_000, format="%d")
    spend_trajectory = st.selectbox("Spend Trajectory", [
        "Growing significantly (>15% YoY)",
        "Growing moderately (5–15% YoY)",
        "Flat (<5% change)",
        "Declining moderately",
        "Declining significantly / Being phased out"
    ])
    geographic_scope = st.multiselect("Geographic Scope", [
        "North America", "Europe", "Asia Pacific", "Latin America",
        "Middle East & Africa", "Global"
    ], default=["North America"])

    st.markdown("---")

    # ── Supplier Landscape ────────────────────────────────────────────────────
    st.markdown("**🏭 Supplier Landscape**")
    num_suppliers = st.number_input("Active Suppliers", min_value=1, value=3, step=1)
    supply_situation = st.selectbox("Supply Situation", [
        "Sole source — no qualified alternatives",
        "Sole source — qualified alternatives exist but not engaged",
        "Preferred supplier with limited competition",
        "2–3 qualified suppliers actively competing",
        "Highly competitive market with many suppliers"
    ])
    last_bid = st.selectbox("Last Competitive Sourcing Event", [
        "Never bid competitively",
        "More than 5 years ago",
        "3–5 years ago",
        "1–3 years ago",
        "Within the last 12 months"
    ])
    incumbent_relationship = st.selectbox("Incumbent Supplier Relationship", [
        "Arms-length / transactional",
        "Cooperative but not strategic",
        "Deeply embedded / difficult to switch",
        "Strategic partner / preferred status",
        "Troubled / under performance management"
    ])

    st.markdown("---")

    # ── Contract & Compliance ─────────────────────────────────────────────────
    st.markdown("**📄 Contract & Compliance**")
    contract_coverage = st.slider("Contract Coverage (%)", 0, 100, 75)
    contract_status = st.selectbox("Contract Status", [
        "No contract in place",
        "Evergreen / auto-renewing",
        "Active — expiring in 12+ months",
        "Active — expiring within 12 months",
        "Active — expiring within 90 days",
        "Recently expired / operating on extension"
    ])
    regulatory_flags = st.multiselect("Applicable Regulatory Frameworks", [
        "FDA / GMP", "ISO 11135", "ISO 11137", "ISO 14644",
        "EMEA Annex 1", "EH&S / OSHA", "REACH / RoHS", "None"
    ])

    st.markdown("---")

    # ── Strategic Context ─────────────────────────────────────────────────────
    st.markdown("**🎯 Strategic Context**")
    primary_objective = st.selectbox("Primary Business Objective", [
        "Cost reduction / savings delivery",
        "Supply assurance / risk mitigation",
        "Quality improvement",
        "Sustainability / ESG compliance",
        "Innovation / supplier development",
        "Compliance / regulatory alignment",
        "Consolidation / simplification"
    ])
    market_maturity = st.selectbox("Market Maturity", ["Emerging", "Growth", "Mature", "Declining"])
    supply_risk = st.select_slider("Perceived Supply Risk", ["Very Low", "Low", "Medium", "High", "Very High"], value="Medium")
    profit_impact = st.select_slider("Profit Impact", ["Very Low", "Low", "Medium", "High", "Very High"], value="Medium")

    procurement_influence = st.selectbox("Procurement's Internal Influence", [
        "Procurement drives — full authority and stakeholder support",
        "Collaborative — procurement leads with strong stakeholder input",
        "Advisory — stakeholders hold final decision authority",
        "Limited — procurement largely executes what business dictates"
    ])

    st.markdown("---")

    # ── Pain Points ───────────────────────────────────────────────────────────
    st.markdown("**⚠️ Current Pain Points**")
    pain_points = st.multiselect("Select All That Apply", [
        "Quality / compliance issues",
        "Delivery / lead time failures",
        "Capacity constraints",
        "Price increases / cost pressure",
        "Poor supplier responsiveness",
        "Invoice / billing disputes",
        "Single point of failure / no backup",
        "Lack of innovation from supplier",
        "Sustainability / ESG concerns",
        "None currently"
    ])

    challenges = st.text_area(
        "Additional Context / Notes",
        placeholder="Anything else the strategy should account for — upcoming audits, internal politics, budget constraints, known market shifts...",
        height=90
    )

    st.markdown("---")

    # ── Spend Data Upload ─────────────────────────────────────────────────────
    st.markdown("**📊 Spend Data (Optional)**")
    st.caption("Upload a CSV to generate spend charts in your deck.")
    spend_file = st.file_uploader("Upload Spend CSV", type=["csv"],
                                  label_visibility="collapsed")
    spend_df = None
    if spend_file:
        try:
            spend_df = pd.read_csv(spend_file)
            supplier_col, amount_col, period_col = detect_columns(spend_df)
            st.success(f"✅ {len(spend_df)} rows loaded")
            st.caption(f"Supplier: `{supplier_col}` · Amount: `{amount_col}`" +
                       (f" · Period: `{period_col}`" if period_col else ""))
        except Exception as e:
            st.error(f"Could not read file: {e}")

    st.markdown("---")
    if not st.session_state.authenticated:
        st.markdown("""
        <div style="background:#1a3a4a; border-radius:8px; padding:1rem; margin-bottom:0.8rem">
            <div style="color:#a8c0cc; font-size:0.78rem; line-height:1.6">
                👁️ <strong style="color:white">Viewing in Demo Mode</strong><br>
                Explore the interface freely. To run a live strategy, contact:<br><br>
                <strong style="color:#00b4d8">Louis Filiano</strong><br>
                Filiano Procurement Consulting<br>
                <a href="mailto:filianoprocurementconsulting@gmail.com" 
                   style="color:#00b4d8; text-decoration:none; font-size:0.78rem">
                   filianoprocurementconsulting@gmail.com
                </a>
            </div>
        </div>
        """, unsafe_allow_html=True)
        pwd = st.text_input("Access Code", type="password", placeholder="Enter access code...",
                            label_visibility="collapsed")
        if st.button("🔓 Unlock", use_container_width=True):
            if pwd == "Birthday-41":
                st.session_state.authenticated = True
                st.rerun()
            else:
                st.error("Incorrect access code.")
        generate_btn = False
    else:
        st.success("✅ Access granted")
        generate_btn = st.button("🚀 Generate Strategy", use_container_width=True)

# ── Session state ─────────────────────────────────────────────────────────────
if "strategy_data" not in st.session_state:
    st.session_state.strategy_data = None
if "raw_output" not in st.session_state:
    st.session_state.raw_output = None

# ── Prompt builder ────────────────────────────────────────────────────────────
def build_prompt(inputs: dict) -> str:
    return f"""You are a world-class Chief Procurement Officer and procurement strategy consultant with 30+ years of experience across medical devices, pharmaceutical, and regulated manufacturing. Generate a comprehensive, board-ready Category Strategy for the following:

CATEGORY DETAILS:
- Category Name: {inputs['category_name']}
- Industry: {inputs['industry']}
- Annual Spend: ${inputs['annual_spend']:,}
- Spend Trajectory: {inputs['spend_trajectory']}
- Active Suppliers: {inputs['num_suppliers']}
- Geographic Scope: {', '.join(inputs['geographic_scope'])}

SUPPLIER DYNAMICS:
- Supply Situation: {inputs['supply_situation']}
- Last Competitive Sourcing Event: {inputs['last_bid']}
- Incumbent Supplier Relationship: {inputs['incumbent_relationship']}

CONTRACT & COMPLIANCE:
- Contract Coverage: {inputs['contract_coverage']}%
- Contract Status: {inputs['contract_status']}
- Regulatory Frameworks: {', '.join(inputs['regulatory_flags']) if inputs['regulatory_flags'] else 'None specified'}

STRATEGIC CONTEXT:
- Primary Business Objective: {inputs['primary_objective']}
- Market Maturity: {inputs['market_maturity']}
- Perceived Supply Risk: {inputs['supply_risk']}
- Profit Impact: {inputs['profit_impact']}
- Procurement's Internal Influence: {inputs['procurement_influence']}

CURRENT PAIN POINTS:
- {', '.join(inputs['pain_points']) if inputs['pain_points'] else 'None identified'}

ADDITIONAL CONTEXT:
- {inputs['challenges'] if inputs['challenges'] else 'None provided'}

Use ALL of the above context to shape every section of the strategy. The supply situation, last bid timing, contract status, primary objective, and pain points should directly influence the sourcing approach, initiatives, risk register, and negotiation levers. Do not produce generic output — make it specific and actionable based on these inputs.

Return ONLY a valid JSON object with NO preamble, NO markdown fences, NO explanation. Use this exact structure:

{{
  "executive_summary": "3-4 sentence strategic summary of the category, its importance, and the recommended approach",
  "category_overview": {{
    "description": "2-3 sentence description of what this category encompasses",
    "spend_classification": "Direct/Indirect/MRO/Services",
    "kraljic_position": "Strategic/Leverage/Bottleneck/Non-Critical",
    "kraljic_rationale": "2-3 sentences explaining the Kraljic positioning",
    "market_outlook": "2-3 sentences on market direction and key trends"
  }},
  "porter_five_forces": {{
    "supplier_power": {{
      "rating": "High/Medium/Low",
      "analysis": "2-3 sentence analysis specific to this category"
    }},
    "buyer_power": {{
      "rating": "High/Medium/Low",
      "analysis": "2-3 sentence analysis"
    }},
    "threat_of_substitutes": {{
      "rating": "High/Medium/Low",
      "analysis": "2-3 sentence analysis"
    }},
    "threat_of_new_entrants": {{
      "rating": "High/Medium/Low",
      "analysis": "2-3 sentence analysis"
    }},
    "competitive_rivalry": {{
      "rating": "High/Medium/Low",
      "analysis": "2-3 sentence analysis"
    }},
    "overall_attractiveness": "Favorable/Neutral/Challenging"
  }},
  "swot": {{
    "strengths": ["strength 1", "strength 2", "strength 3"],
    "weaknesses": ["weakness 1", "weakness 2", "weakness 3"],
    "opportunities": ["opportunity 1", "opportunity 2", "opportunity 3"],
    "threats": ["threat 1", "threat 2", "threat 3"]
  }},
  "sourcing_strategy": {{
    "recommended_approach": "Competitive Bid/Preferred Supplier/Strategic Alliance/Sole Source/Consortium/Other",
    "rationale": "3-4 sentences explaining the recommended sourcing approach",
    "supplier_segmentation": "How suppliers should be tiered and managed",
    "make_vs_buy": "Brief assessment if applicable",
    "negotiation_leverage": ["lever 1", "lever 2", "lever 3"]
  }},
  "key_initiatives": [
    {{
      "initiative": "Initiative title",
      "description": "What it involves and why",
      "timeline": "0-6 months / 6-12 months / 12-24 months / 24-36 months",
      "expected_outcome": "Quantified or qualified benefit",
      "priority": "High/Medium/Low"
    }},
    {{
      "initiative": "Initiative title",
      "description": "What it involves and why",
      "timeline": "0-6 months / 6-12 months / 12-24 months / 24-36 months",
      "expected_outcome": "Quantified or qualified benefit",
      "priority": "High/Medium/Low"
    }},
    {{
      "initiative": "Initiative title",
      "description": "What it involves and why",
      "timeline": "0-6 months / 6-12 months / 12-24 months / 24-36 months",
      "expected_outcome": "Quantified or qualified benefit",
      "priority": "High/Medium/Low"
    }},
    {{
      "initiative": "Initiative title",
      "description": "What it involves and why",
      "timeline": "0-6 months / 6-12 months / 12-24 months / 24-36 months",
      "expected_outcome": "Quantified or qualified benefit",
      "priority": "High/Medium/Low"
    }},
    {{
      "initiative": "Initiative title",
      "description": "What it involves and why",
      "timeline": "0-6 months / 6-12 months / 12-24 months / 24-36 months",
      "expected_outcome": "Quantified or qualified benefit",
      "priority": "High/Medium/Low"
    }}
  ],
  "risk_register": [
    {{
      "risk": "Risk description",
      "likelihood": "High/Medium/Low",
      "impact": "High/Medium/Low",
      "mitigation": "Specific mitigation action"
    }},
    {{
      "risk": "Risk description",
      "likelihood": "High/Medium/Low",
      "impact": "High/Medium/Low",
      "mitigation": "Specific mitigation action"
    }},
    {{
      "risk": "Risk description",
      "likelihood": "High/Medium/Low",
      "impact": "High/Medium/Low",
      "mitigation": "Specific mitigation action"
    }},
    {{
      "risk": "Risk description",
      "likelihood": "High/Medium/Low",
      "impact": "High/Medium/Low",
      "mitigation": "Specific mitigation action"
    }}
  ],
  "kpis": [
    {{"metric": "KPI name", "target": "Specific target", "frequency": "Monthly/Quarterly/Annual"}},
    {{"metric": "KPI name", "target": "Specific target", "frequency": "Monthly/Quarterly/Annual"}},
    {{"metric": "KPI name", "target": "Specific target", "frequency": "Monthly/Quarterly/Annual"}},
    {{"metric": "KPI name", "target": "Specific target", "frequency": "Monthly/Quarterly/Annual"}},
    {{"metric": "KPI name", "target": "Specific target", "frequency": "Monthly/Quarterly/Annual"}}
  ],
  "savings_opportunity": {{
    "estimated_savings_pct": "X-Y% range as a string",
    "estimated_savings_usd": "Dollar amount as a string",
    "primary_lever": "Main savings driver",
    "timeframe": "Expected timeframe to realize savings"
  }}
}}"""


# ── Generate Strategy ─────────────────────────────────────────────────────────
if generate_btn:
    if not category_name:
        st.error("Please enter a Category Name before generating.")
    else:
        inputs = {
            "category_name": category_name,
            "industry": industry,
            "annual_spend": annual_spend,
            "spend_trajectory": spend_trajectory,
            "num_suppliers": num_suppliers,
            "geographic_scope": geographic_scope,
            "supply_situation": supply_situation,
            "last_bid": last_bid,
            "incumbent_relationship": incumbent_relationship,
            "contract_coverage": contract_coverage,
            "contract_status": contract_status,
            "regulatory_flags": regulatory_flags,
            "primary_objective": primary_objective,
            "market_maturity": market_maturity,
            "supply_risk": supply_risk,
            "profit_impact": profit_impact,
            "procurement_influence": procurement_influence,
            "pain_points": pain_points,
            "challenges": challenges,
        }

        with st.spinner("Generating category strategy..."):
            try:
                client = anthropic.Anthropic(api_key=st.secrets["ANTHROPIC_API_KEY"])
                response = client.messages.create(
                    model="claude-opus-4-5",
                    max_tokens=4000,
                    messages=[{"role": "user", "content": build_prompt(inputs)}]
                )
                raw = response.content[0].text.strip()
                # Strip any accidental markdown fences
                if raw.startswith("```"):
                    raw = raw.split("```")[1]
                    if raw.startswith("json"):
                        raw = raw[4:]
                raw = raw.strip()
                data = json.loads(raw)
                st.session_state.strategy_data = data
                st.session_state.raw_output = raw
                st.session_state.inputs = inputs
            except json.JSONDecodeError as e:
                st.error(f"JSON parse error: {e}")
                st.code(raw[:2000])
            except Exception as e:
                st.error(f"Generation error: {e}")

# ── Display Results ───────────────────────────────────────────────────────────
if st.session_state.strategy_data:
    d = st.session_state.strategy_data
    inp = st.session_state.get("inputs", {})

    # ── Top metrics bar ──────────────────────────────────────────────────────
    m1, m2, m3, m4, m5 = st.columns(5)
    with m1:
        st.markdown(f"""<div class="metric-box">
            <div class="value">${inp.get('annual_spend', 0):,.0f}</div>
            <div class="label">Annual Spend</div>
        </div>""", unsafe_allow_html=True)
    with m2:
        st.markdown(f"""<div class="metric-box">
            <div class="value">{inp.get('num_suppliers', '-')}</div>
            <div class="label">Active Suppliers</div>
        </div>""", unsafe_allow_html=True)
    with m3:
        pos = d.get("category_overview", {}).get("kraljic_position", "—")
        st.markdown(f"""<div class="metric-box">
            <div class="value" style="font-size:1.2rem">{pos}</div>
            <div class="label">Kraljic Position</div>
        </div>""", unsafe_allow_html=True)
    with m4:
        savings = d.get("savings_opportunity", {}).get("estimated_savings_pct", "—")
        st.markdown(f"""<div class="metric-box">
            <div class="value" style="font-size:1.3rem">{savings}</div>
            <div class="label">Savings Opportunity</div>
        </div>""", unsafe_allow_html=True)
    with m5:
        attract = d.get("porter_five_forces", {}).get("overall_attractiveness", "—")
        color = {"Favorable": "#68d391", "Neutral": "#f6e05e", "Challenging": "#fc8181"}.get(attract, "white")
        st.markdown(f"""<div class="metric-box">
            <div class="value" style="font-size:1.1rem; color:{color}">{attract}</div>
            <div class="label">Market Attractiveness</div>
        </div>""", unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)

    # ── Tabs ─────────────────────────────────────────────────────────────────
    tabs = st.tabs([
        "📋 Overview",
        "⚡ Five Forces",
        "🔷 SWOT",
        "🎯 Sourcing Strategy",
        "🗺️ Initiatives",
        "⚠️ Risk Register",
        "📊 KPIs",
        "💰 Spend Analysis"
    ])

    # ── Tab 1: Overview ───────────────────────────────────────────────────────
    with tabs[0]:
        ov = d.get("category_overview", {})
        ex = d.get("executive_summary", "")
        
        st.markdown("#### Executive Summary")
        st.markdown(f'<div class="strategy-output">{ex}</div>', unsafe_allow_html=True)

        c1, c2 = st.columns(2)
        with c1:
            st.markdown("#### Category Profile")
            st.markdown(f'<div class="strategy-output">{ov.get("description", "")}</div>', unsafe_allow_html=True)
            st.markdown(f"**Spend Classification:** `{ov.get('spend_classification', '—')}`")
            st.markdown(f"**Kraljic Position:** `{ov.get('kraljic_position', '—')}`")
            st.markdown(f'<div class="strategy-output" style="font-size:0.88rem">{ov.get("kraljic_rationale", "")}</div>', unsafe_allow_html=True)

        with c2:
            st.markdown("#### Market Outlook")
            st.markdown(f'<div class="strategy-output">{ov.get("market_outlook", "")}</div>', unsafe_allow_html=True)
            
            sv = d.get("savings_opportunity", {})
            st.markdown("#### Savings Opportunity")
            st.markdown(f"""
            - **Est. Savings Range:** {sv.get('estimated_savings_pct', '—')}
            - **Est. Dollar Value:** {sv.get('estimated_savings_usd', '—')}
            - **Primary Lever:** {sv.get('primary_lever', '—')}
            - **Timeframe:** {sv.get('timeframe', '—')}
            """)

    # ── Tab 2: Five Forces ────────────────────────────────────────────────────
    with tabs[1]:
        pf = d.get("porter_five_forces", {})
        forces = {
            "🏭 Supplier Power": "supplier_power",
            "🛒 Buyer Power": "buyer_power",
            "🔄 Threat of Substitutes": "threat_of_substitutes",
            "🚪 Threat of New Entrants": "threat_of_new_entrants",
            "⚔️ Competitive Rivalry": "competitive_rivalry",
        }

        col1, col2 = st.columns(2)
        for i, (label, key) in enumerate(forces.items()):
            col = col1 if i % 2 == 0 else col2
            force = pf.get(key, {})
            rating = force.get("rating", "Medium")
            badge_class = {"High": "rating-high", "Medium": "rating-medium", "Low": "rating-low"}.get(rating, "rating-medium")
            with col:
                st.markdown(f"""
                <div class="section-card">
                    <div class="section-label">{label} <span class="rating-badge {badge_class}">{rating}</span></div>
                    <p style="font-size:0.9rem; color:#4a5568; margin:0">{force.get('analysis', '')}</p>
                </div>
                """, unsafe_allow_html=True)

        overall = pf.get("overall_attractiveness", "Neutral")
        color_map = {"Favorable": "#d1fae5", "Neutral": "#fef3c7", "Challenging": "#fde8e8"}
        text_map = {"Favorable": "#065f46", "Neutral": "#92400e", "Challenging": "#c53030"}
        bg = color_map.get(overall, "#fef3c7")
        tc = text_map.get(overall, "#92400e")
        st.markdown(f"""
        <div style="background:{bg}; color:{tc}; padding:1rem 1.5rem; border-radius:10px; text-align:center; font-weight:600; font-size:1.05rem; margin-top:0.5rem">
            Overall Market Attractiveness: {overall}
        </div>
        """, unsafe_allow_html=True)

    # ── Tab 3: SWOT ───────────────────────────────────────────────────────────
    with tabs[2]:
        sw = d.get("swot", {})
        q1, q2 = st.columns(2)
        quadrants = [
            ("💪 Strengths", "strengths", "#d1fae5", "#065f46", q1),
            ("⚠️ Weaknesses", "weaknesses", "#fde8e8", "#c53030", q2),
            ("🚀 Opportunities", "opportunities", "#dbeafe", "#1e40af", q1),
            ("⛈️ Threats", "threats", "#fef3c7", "#92400e", q2),
        ]
        for title, key, bg, tc, col in quadrants:
            items = sw.get(key, [])
            bullets = "".join([f"<li style='margin-bottom:0.4rem'>{item}</li>" for item in items])
            with col:
                st.markdown(f"""
                <div style="background:{bg}; border-radius:10px; padding:1.2rem; margin-bottom:1rem">
                    <div style="font-weight:600; color:{tc}; margin-bottom:0.6rem; font-size:0.95rem">{title}</div>
                    <ul style="margin:0; padding-left:1.2rem; color:#2d3748; font-size:0.88rem">{bullets}</ul>
                </div>
                """, unsafe_allow_html=True)

    # ── Tab 4: Sourcing Strategy ───────────────────────────────────────────────
    with tabs[3]:
        ss = d.get("sourcing_strategy", {})
        st.markdown(f"#### Recommended Approach: `{ss.get('recommended_approach', '—')}`")
        st.markdown(f'<div class="strategy-output">{ss.get("rationale", "")}</div>', unsafe_allow_html=True)

        c1, c2 = st.columns(2)
        with c1:
            st.markdown("**Supplier Segmentation**")
            st.markdown(f'<div class="strategy-output" style="font-size:0.88rem">{ss.get("supplier_segmentation", "")}</div>', unsafe_allow_html=True)
            if ss.get("make_vs_buy"):
                st.markdown("**Make vs. Buy**")
                st.markdown(f'<div class="strategy-output" style="font-size:0.88rem">{ss.get("make_vs_buy", "")}</div>', unsafe_allow_html=True)
        with c2:
            st.markdown("**Negotiation Levers**")
            for lever in ss.get("negotiation_leverage", []):
                st.markdown(f'<span class="tag-pill">🔧 {lever}</span>', unsafe_allow_html=True)

    # ── Tab 5: Initiatives ────────────────────────────────────────────────────
    with tabs[4]:
        initiatives = d.get("key_initiatives", [])
        timeline_order = ["0-6 months", "6-12 months", "12-24 months", "24-36 months"]
        priority_colors = {"High": "#fde8e8", "Medium": "#fef3c7", "Low": "#d1fae5"}
        priority_text = {"High": "#c53030", "Medium": "#92400e", "Low": "#065f46"}

        for ini in sorted(initiatives, key=lambda x: timeline_order.index(x.get("timeline", "12-24 months")) if x.get("timeline") in timeline_order else 99):
            p = ini.get("priority", "Medium")
            bg = priority_colors.get(p, "#fef3c7")
            tc = priority_text.get(p, "#92400e")
            st.markdown(f"""
            <div class="section-card" style="border-left: 4px solid {tc}">
                <div style="display:flex; justify-content:space-between; align-items:flex-start; margin-bottom:0.5rem">
                    <div class="section-label" style="margin:0">📌 {ini.get('initiative', '')}</div>
                    <div>
                        <span class="rating-badge" style="background:{bg}; color:{tc}">{p}</span>
                        <span class="tag-pill">⏱ {ini.get('timeline', '')}</span>
                    </div>
                </div>
                <p style="font-size:0.88rem; color:#4a5568; margin:0.4rem 0">{ini.get('description', '')}</p>
                <p style="font-size:0.85rem; color:#2c5364; font-weight:500; margin:0">✅ {ini.get('expected_outcome', '')}</p>
            </div>
            """, unsafe_allow_html=True)

    # ── Tab 6: Risk Register ──────────────────────────────────────────────────
    with tabs[5]:
        risks = d.get("risk_register", [])
        st.markdown(f"**{len(risks)} identified risks**")

        header_cols = st.columns([3, 1, 1, 3])
        header_cols[0].markdown("**Risk**")
        header_cols[1].markdown("**Likelihood**")
        header_cols[2].markdown("**Impact**")
        header_cols[3].markdown("**Mitigation**")
        st.markdown("---")

        for risk in risks:
            cols = st.columns([3, 1, 1, 3])
            cols[0].write(risk.get("risk", ""))
            lh = risk.get("likelihood", "Medium")
            im = risk.get("impact", "Medium")
            lh_cls = {"High": "rating-high", "Medium": "rating-medium", "Low": "rating-low"}.get(lh, "rating-medium")
            im_cls = {"High": "rating-high", "Medium": "rating-medium", "Low": "rating-low"}.get(im, "rating-medium")
            cols[1].markdown(f'<span class="rating-badge {lh_cls}">{lh}</span>', unsafe_allow_html=True)
            cols[2].markdown(f'<span class="rating-badge {im_cls}">{im}</span>', unsafe_allow_html=True)
            cols[3].write(risk.get("mitigation", ""))
            st.markdown("<hr style='margin:0.4rem 0; border-color:#eee'>", unsafe_allow_html=True)

    # ── Tab 7: KPIs ───────────────────────────────────────────────────────────
    with tabs[6]:
        kpis = d.get("kpis", [])
        for kpi in kpis:
            c1, c2, c3 = st.columns([3, 2, 1])
            c1.markdown(f"**{kpi.get('metric', '')}**")
            c2.markdown(f"`{kpi.get('target', '')}`")
            c3.markdown(f'<span class="tag-pill">{kpi.get("frequency", "")}</span>', unsafe_allow_html=True)
            st.markdown("<hr style='margin:0.5rem 0; border-color:#eee'>", unsafe_allow_html=True)

    # ── Tab 8: Spend Analysis ─────────────────────────────────────────────────
    with tabs[7]:
        if spend_df is not None:
            supplier_col, amount_col, period_col = detect_columns(spend_df)
            total = spend_df[amount_col].sum()
            total_str = f"${total/1_000_000:.2f}M" if total >= 1_000_000 else f"${total/1_000:.0f}K"

            m1, m2, m3 = st.columns(3)
            m1.markdown(f'<div class="metric-box"><div class="value">{total_str}</div><div class="label">Total Spend</div></div>', unsafe_allow_html=True)
            m2.markdown(f'<div class="metric-box"><div class="value">{spend_df[supplier_col].nunique()}</div><div class="label">Suppliers</div></div>', unsafe_allow_html=True)
            top_supplier = spend_df.groupby(supplier_col)[amount_col].sum().idxmax()
            m3.markdown(f'<div class="metric-box"><div class="value" style="font-size:1rem">{top_supplier}</div><div class="label">Top Supplier</div></div>', unsafe_allow_html=True)

            st.markdown("<br>", unsafe_allow_html=True)

            c1, c2 = st.columns(2)
            with c1:
                st.markdown("**Spend by Supplier**")
                bar_img = chart_supplier_bar(spend_df, supplier_col, amount_col)
                st.image(bar_img)
            with c2:
                st.markdown("**Supplier Concentration**")
                pie_img = chart_supplier_pie(spend_df, supplier_col, amount_col)
                st.image(pie_img)

            if period_col and spend_df[period_col].nunique() > 1:
                st.markdown("**Spend Trend**")
                trend_img = chart_spend_trend(spend_df, period_col, amount_col)
                st.image(trend_img, use_container_width=True)

            st.markdown("**Raw Data Preview**")
            st.dataframe(spend_df.head(20), use_container_width=True)
        else:
            st.markdown("""
            <div style="text-align:center; padding:3rem 2rem; color:#718096">
                <div style="font-size:2.5rem; margin-bottom:1rem">📊</div>
                <h4 style="color:#2d3748">No Spend Data Uploaded</h4>
                <p style="max-width:420px; margin:0 auto; font-size:0.9rem">
                    Upload a CSV in the sidebar with columns for Supplier, Amount, and optionally Period to see charts here and in your PowerPoint deck.
                </p>
            </div>
            """, unsafe_allow_html=True)

    # ── Download ──────────────────────────────────────────────────────────────
    st.markdown("---")
    dc1, dc2, dc3 = st.columns([2, 2, 1])
    with dc1:
        pptx_buf = build_pptx(st.session_state.strategy_data,
                              st.session_state.get("inputs", {}),
                              spend_df=spend_df)
        st.download_button(
            "⬇️ Download Strategy (PowerPoint)",
            data=pptx_buf,
            file_name=f"category_strategy_{category_name.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    with dc2:
        raw_json = json.dumps(st.session_state.strategy_data, indent=2)
        st.download_button(
            "⬇️ Download Strategy (JSON)",
            data=raw_json,
            file_name=f"category_strategy_{category_name.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.json",
            mime="application/json"
        )
    with dc3:
        st.caption(f"Generated {datetime.now().strftime('%B %d, %Y at %I:%M %p')}")

else:
    # ── Empty state ───────────────────────────────────────────────────────────
    st.markdown("""
    <div style="text-align:center; padding:4rem 2rem; color:#718096">
        <div style="font-size:3rem; margin-bottom:1rem">🎯</div>
        <h3 style="font-family:'DM Serif Display', serif; color:#2d3748">Ready to Build Your Category Strategy</h3>
        <p style="max-width:500px; margin:0 auto; font-size:0.95rem; line-height:1.7">
            Fill in the category details in the sidebar and click <strong>Generate Strategy</strong>. 
            The AI will produce a complete, board-ready strategy including Porter's Five Forces, 
            Kraljic positioning, SWOT, sourcing approach, initiatives, risk register, and KPIs.
        </p>
    </div>
    """, unsafe_allow_html=True)